import streamlit as st
import requests
import json
import os
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Streamlit UI
st.title("Presentation Generation System")

# API base URL (update if running on a different host/port)
API_URL = "http://localhost:8000"

# File upload
uploaded_file = st.file_uploader("Upload a PDF document", type=["pdf"])
if uploaded_file:
    with st.spinner("Uploading PDF..."):
        files = {"file": (uploaded_file.name, uploaded_file, "application/pdf")}
        response = requests.post(f"{API_URL}/upload_pdf", files=files)
        if response.status_code == 200:
            data = response.json()
            session_id = data["session_id"]
            st.session_state["session_id"] = session_id
            st.success("PDF uploaded successfully!")
        else:
            st.error("Failed to upload PDF.")

# Query input and slide generation
if "session_id" in st.session_state:
    query = st.text_input("Enter your presentation topic:", value="Generate a presentation about the Assessment List for Trustworthy AI (ALTAI) for a general audience.")
    if st.button("Generate Slides"):
        with st.spinner("Generating slides... This may take a few minutes."):
            payload = {"query": query, "session_id": st.session_state["session_id"]}
            response = requests.post(f"{API_URL}/generate_slides", json=payload)
            if response.status_code == 200:
                data = response.json()
                slides = data["slides"]
                st.session_state["slides"] = slides
                st.success("Slides generated successfully!")
            else:
                st.error(f"Error: {response.json()['detail']}")

# Function to convert slides to PPT
def slides_to_ppt(slides):
    """Convert a list of slides (JSON) to a PowerPoint presentation."""
    prs = Presentation()

    # Define slide layout (Title and Content)
    slide_layout = prs.slide_layouts[1]  # 1 is the "Title and Content" layout

    for slide_data in slides:
        # Add a new slide
        slide = prs.slides.add_slide(slide_layout)

        # Set the title
        title_shape = slide.shapes.title
        title_shape.text = slide_data["title"]

        # Set the content (body)
        body_shape = slide.placeholders[1]  # Placeholder for content
        text_frame = body_shape.text_frame
        text_frame.text = slide_data["content"]

        # Optional: Format the content
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.size = Pt(18)

    # Save the presentation to a BytesIO buffer
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

# Display and download slides
if "slides" in st.session_state:
    st.header("Generated Slides")
    for i, slide in enumerate(st.session_state["slides"], 1):
        st.subheader(f"Slide {i}: {slide['title']}")
        st.write(slide["content"])
        st.markdown("---")

    # Download option for JSON
    slides_json = json.dumps(st.session_state["slides"], indent=4)
    st.download_button(
        label="Download Slides as JSON",
        data=slides_json,
        file_name="presentation.json",
        mime="application/json"
    )

    # Download option for PPT
    ppt_buffer = slides_to_ppt(st.session_state["slides"])
    st.download_button(
        label="Download Slides as PPT",
        data=ppt_buffer,
        file_name="presentation.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )